"""
PPT 생성 모듈 - 6가지 슬라이드 유형 생성
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from typing import Dict, List, Optional, Any
import pandas as pd

from . import config
from .dropbox_client import get_dropbox_client
from .sheets_client import get_sheets_client


class PPTGenerator:
    """심의자료 PPT 생성기"""
    
    def __init__(self):
        self.dropbox = get_dropbox_client()
        self.sheets = get_sheets_client()
        self.ppt: Optional[Presentation] = None
        self.materials_sizes: Dict[str, List[str]] = {}
        self.materials: List[str] = []
        self.text_assets: Dict[str, Any] = {}
        self.df_obj_result: Optional[pd.DataFrame] = None
    
    def generate(self, keyword: str, progress_callback=None) -> BytesIO:
        """
        PPT 생성 메인 함수
        Args:
            keyword: 소재 필터링 키워드 (예: "usp-dm-1st")
            progress_callback: 진행 상태 콜백 함수
        Returns:
            PPT 파일의 BytesIO 객체
        """
        # 1. PPT 템플릿 로드 (BytesIO에서)
        self._load_template()
        
        # 2. 소재 목록 조회
        self.materials_sizes = self.dropbox.get_materials_list(keyword)
        self.materials = list(self.materials_sizes.keys())
        
        if not self.materials:
            raise ValueError(f"키워드 '{keyword}'에 해당하는 소재가 없습니다.")
        
        # 3. 텍스트 에셋 로드
        self.text_assets = self.sheets.get_text_assets(self.materials[0])
        self.df_obj_result = self.sheets.get_object_assets(self.materials)
        
        # 4. 슬라이드 생성 (6가지 유형)
        self._create_all_slides(progress_callback)
        
        # 5. BytesIO로 저장
        ppt_buffer = BytesIO()
        self.ppt.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
    
    def generate_with_materials(self, selected_materials: List[str], progress_callback=None) -> BytesIO:
        """
        선택된 소재로 PPT 생성
        Args:
            selected_materials: 선택된 소재명 리스트
            progress_callback: 진행 상태 콜백 함수
        Returns:
            PPT 파일의 BytesIO 객체
        """
        # 1. PPT 템플릿 로드
        self._load_template()
        
        # 2. 선택된 소재들의 사이즈 정보 조회
        all_materials = self.dropbox.get_materials_list(None)
        self.materials_sizes = {m: all_materials[m] for m in selected_materials if m in all_materials}
        self.materials = list(self.materials_sizes.keys())
        
        if not self.materials:
            raise ValueError("선택된 소재를 찾을 수 없습니다.")
        
        # 3. 텍스트 에셋 로드
        self.text_assets = self.sheets.get_text_assets(self.materials[0])
        self.df_obj_result = self.sheets.get_object_assets(self.materials)
        
        # 4. 슬라이드 생성 (6가지 유형)
        self._create_all_slides(progress_callback)
        
        # 5. BytesIO로 저장
        ppt_buffer = BytesIO()
        self.ppt.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
    
    def generate_with_progress(self, selected_materials: List[str], progress_callback=None) -> BytesIO:
        """
        진행상황 콜백 포함 PPT 생성 (상세 단계별)
        Args:
            selected_materials: 선택된 소재명 리스트
            progress_callback: 콜백(step, current, total, detail)
        """
        def notify(step: str, current: int, total: int, detail: str = ""):
            if progress_callback:
                progress_callback(step, current, total, detail)
        
        # 1. 초기화
        notify("초기화", 0, 100, "템플릿 로드 중...")
        self._load_template()
        
        # 2. 소재 정보 조회
        notify("소재 조회", 5, 100, f"{len(selected_materials)}개 소재 정보 확인 중...")
        all_materials = self.dropbox.get_materials_list(None)
        self.materials_sizes = {m: all_materials[m] for m in selected_materials if m in all_materials}
        self.materials = list(self.materials_sizes.keys())
        
        if not self.materials:
            raise ValueError("선택된 소재를 찾을 수 없습니다.")
        
        # 3. 이미지 프리로드 (병렬)
        notify("이미지 다운로드", 10, 100, "이미지 병렬 다운로드 시작...")
        
        # 필요한 모든 사이즈 목록
        all_sizes = [
            '640x100', '970x250', '160x600',
            '1200x628', '1200x1200', '1200x1500',
            '1080x1080', '1200x1200_toss',
            '315x258', '342x228', '112x112',
            '200x200_toss', '1200x1200_당근',
            '1200x627_CTAx'
        ]
        
        def img_progress(current, total, msg):
            percent = 10 + int((current / total) * 40)  # 10% ~ 50%
            notify("이미지 다운로드", percent, 100, msg)
        
        self.dropbox.preload_images(self.materials, all_sizes, img_progress)
        
        # 4. 텍스트 에셋 로드
        notify("텍스트 로드", 52, 100, "Google Sheets에서 텍스트 에셋 로드 중...")
        self.text_assets = self.sheets.get_text_assets(self.materials[0])
        self.df_obj_result = self.sheets.get_object_assets(self.materials)
        
        # 5. 슬라이드 생성
        slides_funcs = [
            ("배너형 슬라이드", self._first_create_slides),
            ("정사각/세로형 슬라이드", self._second_create_slides),
            ("구글 텍스트에셋", self._third_create_slides),
            ("META/토스 모먼트탭", self._fourth_create_slides),
            ("오브젝트형", self._fifth_create_slides),
            ("버즈빌/스페셜DA/GFA", self._sixth_create_slides),
        ]
        
        for i, (name, func) in enumerate(slides_funcs):
            percent = 55 + int((i / len(slides_funcs)) * 40)  # 55% ~ 95%
            notify("슬라이드 생성", percent, 100, f"📄 {name} 생성 중...")
            func()
        
        # 6. 저장
        notify("완료", 98, 100, "PPT 파일 저장 중...")
        ppt_buffer = BytesIO()
        self.ppt.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        notify("완료", 100, 100, "✅ 생성 완료!")
        return ppt_buffer
    
    def _create_all_slides(self, progress_callback=None):
        """모든 슬라이드 유형 생성"""
        slides_funcs = [
            ("배너형 슬라이드", self._first_create_slides),
            ("정사각/세로형 슬라이드", self._second_create_slides),
            ("구글 텍스트에셋", self._third_create_slides),
            ("META/토스 모먼트탭", self._fourth_create_slides),
            ("오브젝트형", self._fifth_create_slides),
            ("버즈빌/스페셜DA/GFA", self._sixth_create_slides),
        ]
        
        total = len(slides_funcs)
        for i, (name, func) in enumerate(slides_funcs):
            if progress_callback:
                progress = int(20 + (i / total) * 70)
                progress_callback(name, progress, 100, f"📄 {name}")
            func()
        
        if progress_callback:
            progress_callback("슬라이드 생성 완료", 95, 100, "저장 중...")
    
    def _load_template(self):
        """PPT 템플릿 로드"""
        # 로컬 템플릿 파일 사용 (Dropbox에서 다운로드도 가능)
        if config.TEMPLATE_PATH.exists():
            self.ppt = Presentation(str(config.TEMPLATE_PATH))
        else:
            # 빈 프레젠테이션 생성 (폴백)
            self.ppt = Presentation()
    
    def _add_image_from_dropbox(self, slide, material: str, size: str, 
                                 left: float, top: float, width: float, height: float) -> bool:
        """Dropbox에서 이미지 다운로드 후 슬라이드에 추가"""
        try:
            img_bytes = self.dropbox.download_image(material, size)
            if img_bytes:
                slide.shapes.add_picture(img_bytes, Cm(left), Cm(top), 
                                        width=Cm(width), height=Cm(height))
                return True
        except Exception as e:
            print(f"Error adding image {material}/{size}: {e}")
        return False
    
    # ===== 슬라이드 유형 1: 배너형 =====
    def _first_create_slides(self):
        """첫번째 유형: 640x100, 970x250, 160x600"""
        for i in range(0, len(self.materials), 2):
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[0])
            
            # 첫번째 소재
            material = self.materials[i]
            sizes = self.materials_sizes[material]
            
            # 640x100
            if '640x100' in sizes:
                self._add_image_from_dropbox(slide, material, '640x100', 0, 0.7, 18.54, 2.9)
            
            # 970x250
            if '970x250' in sizes:
                self._add_image_from_dropbox(slide, material, '970x250', 0, 7.38, 18.59, 4.79)
            
            # 160x600
            if '160x600' in sizes:
                self._add_image_from_dropbox(slide, material, '160x600', 18.61, 0.7, 4.48, 16.8)
            
            # 두번째 소재
            if i + 1 < len(self.materials):
                material2 = self.materials[i + 1]
                sizes2 = self.materials_sizes[material2]
                
                if '640x100' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '640x100', 0, 3.94, 18.54, 2.9)
                
                if '970x250' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '970x250', 0, 12.71, 18.59, 4.79)
                
                if '160x600' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '160x600', 23.09, 0.7, 4.48, 16.8)
    
    # ===== 슬라이드 유형 2: 정사각/세로형 =====
    def _second_create_slides(self):
        """두번째 유형: 1200x628, 1200x1200, 1200x1500"""
        for i in range(0, len(self.materials), 2):
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[1])
            
            # 첫번째 소재
            material = self.materials[i]
            sizes = self.materials_sizes[material]
            
            if '1200x628' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x628', 0, 0.18, 13.86, 7.25)
            
            if '1200x1200' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1200', 0, 7.62, 7.15, 7.15)
            
            if '1200x1500' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1500', 7.27, 7.62, 6.4, 8)
            
            # 두번째 소재
            if i + 1 < len(self.materials):
                material2 = self.materials[i + 1]
                sizes2 = self.materials_sizes[material2]
                
                if '1200x628' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x628', 13.65, 0.18, 13.86, 7.25)
                
                if '1200x1200' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x1200', 13.65, 7.62, 7.15, 7.15)
                
                if '1200x1500' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x1500', 21.12, 7.62, 6.4, 8)
    
    # ===== 슬라이드 유형 3: 구글 텍스트에셋 =====
    def _third_create_slides(self):
        """세번째 유형: 이미지 + 구글 텍스트에셋 테이블"""
        google_list = self.text_assets.get('google_range_list', [])
        if not google_list:
            return
        
        for i in range(0, len(self.materials), 2):
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[2])
            
            # 첫번째 소재
            material = self.materials[i]
            sizes = self.materials_sizes[material]
            
            if '1200x628' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x628', 19.15, 8.21, 8.24, 4.31)
            
            if '1200x1200' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1200', 0, 0.18, 7.15, 7.15)
            
            if '1200x1500' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1500', 7.27, 0.18, 6.4, 8)
            
            # 텍스트 테이블 추가
            self._add_google_text_tables(slide, google_list)
            
            # 두번째 소재
            if i + 1 < len(self.materials):
                material2 = self.materials[i + 1]
                sizes2 = self.materials_sizes[material2]
                
                if '1200x628' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x628', 19.15, 12.53, 8.24, 4.31)
                
                if '1200x1200' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x1200', 13.75, 0.18, 7.15, 7.15)
                
                if '1200x1500' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x1500', 21.12, 0.18, 6.4, 8)
    
    def _add_google_text_tables(self, slide, google_list: List[str]):
        """구글 텍스트에셋 테이블 추가"""
        table_configs = [
            (8.58, 0, 5),    # 표1: 0-4, 5-9
            (11.59, 10, 15), # 표2: 10-14, 15-19
            (14.57, 20, 25), # 표3: 20-24, 25-29
        ]
        
        for top_cm, title_start, desc_start in table_configs:
            # 제목 테이블
            self._create_text_table(slide, 5, 1, 1.68, top_cm, 5.37, 2.93,
                                   google_list[title_start:title_start+5])
            # 설명 테이블
            self._create_text_table(slide, 5, 1, 8.72, top_cm, 10.22, 2.93,
                                   google_list[desc_start:desc_start+5])
    
    def _create_text_table(self, slide, rows: int, cols: int, 
                           left: float, top: float, width: float, height: float,
                           texts: List[str]):
        """테이블 생성 및 텍스트 채우기"""
        table = slide.shapes.add_table(rows, cols, Cm(left), Cm(top), 
                                       Cm(width), Cm(height)).table
        table.columns[0].width = Cm(width)
        
        for i, text in enumerate(texts):
            if i < rows:
                cell = table.cell(i, 0)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.text = text if text else ''
    
    # ===== 슬라이드 유형 4: META/토스 모먼트탭 =====
    def _fourth_create_slides(self):
        """네번째 유형: META, 토스 모먼트탭"""
        meta_list = self.text_assets.get('meta_range_list', [])
        meta_caution = self.text_assets.get('meta_caution', '')
        
        for i in range(0, len(self.materials), 2):
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[3])
            
            # 첫번째 소재
            material = self.materials[i]
            sizes = self.materials_sizes[material]
            
            # 1080x1080
            if '1080x1080' in sizes:
                self._add_image_from_dropbox(slide, material, '1080x1080', 0.64, 9.49, 6.63, 6.63)
            
            # META 텍스트
            self._add_meta_text(slide, meta_list, meta_caution, 0.64)
            
            # 1200x1200_toss
            if '1200x1200_toss' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1200_toss', 15.14, 0.5, 5.41, 5.41)
            
            # 토스 모먼트탭 텍스트
            self._add_toss_moment_text(slide, material, 15.14)
            
            # 두번째 소재
            if i + 1 < len(self.materials):
                material2 = self.materials[i + 1]
                sizes2 = self.materials_sizes[material2]
                
                if '1080x1080' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1080x1080', 7.81, 9.49, 6.63, 6.63)
                
                self._add_meta_text(slide, meta_list, meta_caution, 7.81)
                
                if '1200x1200_toss' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x1200_toss', 21.47, 0.5, 5.41, 5.41)
                
                self._add_toss_moment_text(slide, material2, 21.47)
    
    def _add_meta_text(self, slide, meta_list: List[str], meta_caution: str, left: float):
        """META 텍스트 추가"""
        # 본문
        if len(meta_list) >= 7:
            body_text = '\n'.join(meta_list[1:7])
            self._add_textbox(slide, left, 0.8, 6.52, 4, body_text, 7.5)
        
        # 유의문구
        if meta_caution:
            caution_wrapped = self._add_newlines(meta_caution, 28)
            self._add_textbox(slide, left, 2.8, 6.52, 4, caution_wrapped, 7.5)
        
        # 제목
        if meta_list:
            self._add_textbox(slide, left, 16.14, 6.52, 2, meta_list[0], 9)
    
    def _add_toss_moment_text(self, slide, material: str, left: float):
        """토스 모먼트탭 텍스트 추가"""
        try:
            text1 = self.df_obj_result.loc[material, "토스_모먼트탭_메인문구1"]
            text2 = self.df_obj_result.loc[material, "토스_모먼트탭_메인문구2"]
            text3 = self.df_obj_result.loc[material, "토스_모먼트탭_보조문구"]
            full_text = f"{text1}\n{text2}\n{text3}"
            self._add_textbox(slide, left, 6.18, 5.5, 5.88, full_text, 10, 
                            font_color=RGBColor(255, 255, 255))
        except:
            pass
    
    # ===== 슬라이드 유형 5: 오브젝트형 =====
    def _fifth_create_slides(self):
        """다섯번째 유형: 카카오/네이버/토스/당근"""
        for material in self.materials:
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[4])
            sizes = self.materials_sizes[material]
            
            # 카카오 비즈보드
            self._add_textbox(slide, 1.89, 4.25, 7, 0.92,
                            self._get_obj_value(material, "카카오_비즈보드_메인카피"), 14)
            self._add_textbox(slide, 1.89, 5.0, 7, 0.92,
                            self._get_obj_value(material, "카카오_비즈보드_서브카피"), 12)
            
            if '315x258' in sizes:
                self._add_image_from_dropbox(slide, material, '315x258', 9.9, 4, 3, 2.3)
            
            # 몰로코 비즈보드
            self._add_textbox(slide, 14.7, 4.25, 7, 0.92,
                            self._get_obj_value(material, "카카오_비즈보드(몰로코,애피어)_메인카피"), 14)
            
            if '315x258' in sizes:
                self._add_image_from_dropbox(slide, material, '315x258', 22.16, 4.25, 3, 2.3)
            
            # 네이버 네이티브
            self._add_textbox(slide, 1.77, 7.45, 7, 0.92,
                            self._get_obj_value(material, "네이버GFA_네이티브_광고문구"), 12)
            self._add_textbox(slide, 6.03, 8.53, 7, 0.92,
                            self._get_obj_value(material, "네이버GFA_네이티브_설명문구1"), 10)
            self._add_textbox(slide, 6.03, 9.03, 7, 0.92,
                            self._get_obj_value(material, "네이버GFA_네이티브_설명문구2"), 10)
            self._add_textbox(slide, 6.03, 9.53, 7, 0.92,
                            self._get_obj_value(material, "네이버GFA_네이티브_설명문구3"), 10)
            
            if '342x228' in sizes:
                self._add_image_from_dropbox(slide, material, '342x228', 1.93, 9, 4, 2.62)
            
            # 네이버 커뮤니케이션애드
            text1 = self._add_newlines(self._get_obj_value(material, "네이버GFA_커뮤니케이션애드_광고문구1"), 23)
            self._add_textbox(slide, 12.5, 7.57, 7, 0.92, text1, 9.5)
            
            if '112x112' in sizes:
                self._add_image_from_dropbox(slide, material, '112x112', 20.21, 8.27, 1.77, 1.77)
            
            text2 = self._add_newlines(self._get_obj_value(material, "네이버GFA_커뮤니케이션애드_광고문구2"), 23)
            self._add_textbox(slide, 13.79, 10.95, 8, 0.92, text2, 9.5)
            
            # 토스 혜택탭
            if '200x200_toss' in sizes:
                self._add_image_from_dropbox(slide, material, '200x200_toss', 1.87, 14.33, 1.27, 1.27)
            
            self._add_textbox(slide, 3.15, 13.34, 7, 0.92,
                            self._get_obj_value(material, "토스_혜택탭_메인문구"), 13)
            sub_text = self._get_obj_value(material, "토스_혜택탭_보조문구") + " AD"
            self._add_textbox(slide, 3.15, 13.94, 7, 0.92, sub_text, 11.5)
            
            # 당근 네이티브
            if '1200x1200_당근' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1200_당근', 14.0, 14.46, 2.9, 2.9)
            
            title = self._add_newlines(self._get_obj_value(material, "당근_당근네이티브_광고 제목"), 20)
            self._add_textbox(slide, 17.28, 13.46, 7, 1.05, title, 13.5)
            
            caution = self._add_newlines(self._get_obj_value(material, "당근_당근네이티브_심의필 문구"), 20)
            self._add_textbox(slide, 17.28, 15.29, 7, 1.05, caution, 8)
    
    # ===== 슬라이드 유형 6: 버즈빌/스페셜DA/GFA홈피드 =====
    def _sixth_create_slides(self):
        """여섯번째 유형: 버즈빌, 스페셜DA, GFA홈피드"""
        for i in range(0, len(self.materials), 2):
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[9])
            
            # 버즈빌 텍스트
            self._add_textbox(slide, 2.72, 2.36, 6, 2,
                            self._get_obj_value(self.materials[0], "버즈빌_카카오금융_광고 제목"), 9)
            
            # 첫번째 소재
            material = self.materials[i]
            sizes = self.materials_sizes[material]
            
            if '1200x627_CTAx' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x627_CTAx', 2.72, 4.98, 10.34, 5.4)
            
            # 스페셜DA 텍스트
            self._add_textbox(slide, 16.64, 1.75, 4.5, 2,
                            self._get_obj_value(material, "카카오_비즈보드_메인카피"), 7.5)
            self._add_textbox(slide, 16.64, 2.4, 4.5, 2,
                            self._get_obj_value(material, "카카오_비즈보드_서브카피"), 7.5)
            
            if '315x258' in sizes:
                self._add_image_from_dropbox(slide, material, '315x258', 21.57, 2.04, 2.43, 2)
            
            # GFA 홈피드
            self._add_textbox(slide, 13.76, 6.8, 7, 2,
                            self._get_obj_value(material, "네이버GFA_네이티브_광고문구"), 7)
            
            desc_text = f"{self._get_obj_value(material, '네이버GFA_네이티브_설명문구1')} {self._get_obj_value(material, '네이버GFA_네이티브_설명문구2')}\n {self._get_obj_value(material, '네이버GFA_네이티브_설명문구3')}"
            self._add_textbox(slide, 13.76, 7.15, 5.94, 2, desc_text, 7.5)
            
            if '1200x1200' in sizes:
                self._add_image_from_dropbox(slide, material, '1200x1200', 13.61, 9.08, 6.1, 6.1)
            
            # 두번째 소재
            if i + 1 < len(self.materials):
                material2 = self.materials[i + 1]
                sizes2 = self.materials_sizes[material2]
                
                if '1200x627_CTAx' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x627_CTAx', 2.72, 10.63, 10.34, 5.4)
                
                self._add_textbox(slide, 16.64, 3.93, 4.5, 2,
                                self._get_obj_value(material2, "카카오_비즈보드_메인카피"), 7.5)
                self._add_textbox(slide, 16.64, 4.58, 4.5, 2,
                                self._get_obj_value(material2, "카카오_비즈보드_서브카피"), 7.5)
                
                if '315x258' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '315x258', 21.57, 4.23, 2.43, 2)
                
                self._add_textbox(slide, 20.2, 6.8, 7, 2,
                                self._get_obj_value(material2, "네이버GFA_네이티브_광고문구"), 7)
                
                desc_text2 = f"{self._get_obj_value(material2, '네이버GFA_네이티브_설명문구1')} {self._get_obj_value(material2, '네이버GFA_네이티브_설명문구2')}\n {self._get_obj_value(material2, '네이버GFA_네이티브_설명문구3')}"
                self._add_textbox(slide, 20.2, 7.15, 5.94, 2, desc_text2, 7.5)
                
                if '1200x1200' in sizes2:
                    self._add_image_from_dropbox(slide, material2, '1200x1200', 20.2, 9.08, 6.1, 6.1)
    
    # ===== 헬퍼 함수들 =====
    def _add_textbox(self, slide, left: float, top: float, width: float, height: float,
                     text: str, font_size: float, font_color: RGBColor = RGBColor(0, 0, 0)):
        """텍스트박스 추가"""
        txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.text = text if text else ''
        p.font.name = "Malgun Gothic"
    
    def _add_newlines(self, text: str, max_chars: int) -> str:
        """지정된 글자수에서 줄바꿈 추가"""
        if not text or len(text) <= max_chars:
            return text
        return text[:max_chars] + '\n' + text[max_chars:]
    
    def _get_obj_value(self, material: str, column: str) -> str:
        """DataFrame에서 값 가져오기"""
        try:
            return str(self.df_obj_result.loc[material, column])
        except:
            return ''


def generate_ppt(keyword: str) -> BytesIO:
    """PPT 생성 편의 함수"""
    generator = PPTGenerator()
    return generator.generate(keyword)
